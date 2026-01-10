#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
スピーカーノートからPowerPointスライドを生成するスクリプト
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_speaker_notes(file_path):
    """スピーカーノートファイルを解析してスライド情報を抽出"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    # まず、すべてのスライドセクションを検出
    slide_sections = re.split(r'## スライド(\d+):\s*(.+?)\n', content)
    
    # 最初の要素はヘッダーなのでスキップ、その後は交互に番号、タイトル、内容
    for i in range(1, len(slide_sections), 3):
        if i + 2 < len(slide_sections):
            slide_num = int(slide_sections[i])
            title = slide_sections[i + 1].strip()
            section_content = slide_sections[i + 2]
            
            # 「**話す内容:**」の後の内容を抽出
            content_match = re.search(r'\*\*話す内容:\*\*\n((?:- .+\n?)+)', section_content)
            if content_match:
                content_lines = content_match.group(1)
                
                # 箇条書きの内容を抽出（簡潔なポイントのみ）
                points = []
                for line in content_lines.split('\n'):
                    line = line.strip()
                    if line.startswith('- 「') and '」' in line:
                        # 「」の間のテキストを抽出
                        text_match = re.search(r'「(.+?)」', line)
                        if text_match:
                            text = text_match.group(1)
                            # スライド用に簡潔にする（主要なポイントのみ）
                            # 長すぎる文や詳細な説明は除外
                            if len(text) <= 100:
                                # 重要なポイントのみを抽出
                                # 体験談や詳細説明は除外
                                if not any(exclude in text for exclude in ['僕も', '最初は', '続けていくうちに', '実はそれが']):
                                    # キーワードを含む重要なポイント、または最初の数スライド
                                    if any(keyword in text for keyword in ['原則', '法則', '重要', '大事', '確認', '方法', '解決', 'まとめ', '領域', '市場', '競合', '需要', '差別化', 'CORE', '成長', '経験', '権威', '再現', '情報統合', '定量', 'ライバル', 'フレームワーク']):
                                        if len(text) <= 60:  # さらに短く
                                            points.append(text)
                                    elif slide_num <= 3:  # 最初の数スライドは多めに
                                        if len(text) <= 60:
                                            points.append(text)
                
                # ポイントが多すぎる場合は最初の3-4個に絞る
                if len(points) > 4:
                    points = points[:4]
                
                slides.append({
                    'number': slide_num,
                    'title': title,
                    'points': points
                })
    
    # スライド番号でソート
    slides.sort(key=lambda x: x['number'])
    
    return slides

def create_presentation(slides, output_path):
    """PowerPointプレゼンテーションを作成"""
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # タイトルスライドを最初に追加
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "ビジネス領域選定実践プログラム"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "自分が本当に取り組むべき大きな領域を見つける"
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    
    # 各スライドを作成
    for slide_data in slides:
        # タイトルとコンテンツのレイアウトを使用
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title_shape = slide.shapes.title
        title_shape.text = slide_data['title']
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # コンテンツ
        if slide_data['points']:
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.word_wrap = True
            tf.clear()  # デフォルトの段落をクリア
            
            for i, point in enumerate(slide_data['points']):
                p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(24)
                p.space_after = Pt(12)
                p.level = 0
        else:
            # ポイントがない場合は説明文を追加
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.text = "詳細はスピーカーノートを参照してください"
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.italic = True
    
    prs.save(output_path)
    print(f"PowerPointファイルを生成しました: {output_path}")
    print(f"合計 {len(slides) + 1} スライド（タイトル含む）")

def main():
    input_file = "オリジナル講座/第一章 1-3 ビジネス領域選定実践プログラム_スピーカーノート.md"
    output_file = "オリジナル講座/第一章 1-3 ビジネス領域選定実践プログラム.pptx"
    
    print("スピーカーノートを解析中...")
    slides = parse_speaker_notes(input_file)
    print(f"{len(slides)}個のスライドを検出しました")
    
    print("PowerPointファイルを生成中...")
    create_presentation(slides, output_file)
    print("完了しました！")

if __name__ == "__main__":
    main()

